import os


def extract_text_via_vision(file_path, course_code="", batch_size=15):
    import fitz
    import tempfile
    import time
    from google import genai
    from google.genai import types
    from django.conf import settings

    api_key = getattr(settings, "GEMINI_API_KEY_EXTRACTION", None)
    if not api_key:
        raise ValueError("GEMINI_API_KEY_EXTRACTION is missing from settings.")

    client = genai.Client(api_key=api_key)

    doc = fitz.open(file_path)
    total_pages = len(doc)
    doc.close()

    print(f"[{course_code}] {total_pages} pages total — transcribing in batches of {batch_size}.")

    transcript_parts = []

    for batch_start in range(0, total_pages, batch_size):
        batch_end = min(batch_start + batch_size, total_pages) - 1
        print(f"[{course_code}] Batch pages {batch_start + 1}-{batch_end + 1}...")

        doc = fitz.open(file_path)
        batch_doc = fitz.open()
        batch_doc.insert_pdf(doc, from_page=batch_start, to_page=batch_end)
        doc.close()

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_path = temp_pdf.name
        batch_doc.save(temp_path)
        batch_doc.close()

        gemini_file = None
        batch_text = ""

        try:
            for attempt in range(3):
                try:
                    with open(temp_path, "rb") as f:
                        gemini_file = client.files.upload(file=f, config={"mime_type": "application/pdf"})
                    break
                except Exception as e:
                    print(f"[{course_code}] Batch upload attempt {attempt + 1} failed: {e}")
                    if attempt < 2:
                        time.sleep(3)
                    else:
                        raise

            for model in ["gemini-3.6-flash", "gemini-3.7-flash"]:
                try:
                    response = client.models.generate_content(
                        model=model,
                        contents=[
                            gemini_file,
                            f"Transcribe all text and handwriting in this document exactly as written, "
                            f"page by page, including formulas, headings, and diagram labels. "
                            f"Preserve structure with line breaks. Mark each page as '--- Page N ---', "
                            f"using the ACTUAL page number in the full document (this batch starts at "
                            f"page {batch_start + 1}). For handwritten formulas, use plain text math "
                            f"notation. Do not summarize or explain — transcribe only.",
                        ],
                        config=types.GenerateContentConfig(max_output_tokens=8000),
                    )
                    batch_text = response.text.strip()
                    print(f"[{course_code}] Batch done using {model} ({len(batch_text)} chars).")
                    break
                except Exception as e:
                    print(f"[{course_code}] Model {model} failed on batch: {e}. Trying next...")

            if not batch_text:
                print(f"[{course_code}] WARNING: batch pages {batch_start + 1}-{batch_end + 1} produced no text.")

        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            if gemini_file:
                try:
                    client.files.delete(name=gemini_file.name)
                except Exception as e:
                    print(f"[{course_code}] Could not delete remote batch file: {e}")

        transcript_parts.append(batch_text)

    full_transcript = "\n\n".join(transcript_parts)
    if not full_transcript.strip():
        raise RuntimeError("All batches failed to transcribe the document.")

    print(f"[{course_code}] Full transcription complete: {len(full_transcript)} chars across {total_pages} pages.")
    return full_transcript


def extract_text_from_slide(file_path, course_code=""):
    path = str(file_path)

    if path.endswith(".pdf"):
        return extract_text_via_vision(path, course_code=course_code)

    elif path.endswith(".docx"):
        import docx
        doc = docx.Document(path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text.strip()

    elif path.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(path)
        extracted_chunks = []

        def extract_from_shape(shape):
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        extracted_chunks.append(para_text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            extracted_chunks.append(cell_text)
            elif shape.shape_type == 6:
                for sub_shape in shape.shapes:
                    extract_from_shape(sub_shape)

        for slide in prs.slides:
            for shape in slide.shapes:
                extract_from_shape(shape)

        return "\n".join(extracted_chunks).strip()

    else:
        return ""